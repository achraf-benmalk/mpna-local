#!/usr/bin/env python3
"""Generate the HPL presentation slides (Part 2: Execution, Results, Analysis).
Styled to match Presentation_Nguyen2014_5min_V2.pptx template."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

IMG_DIR = os.path.join(os.path.dirname(__file__), '..', 'latex report', 'images')
OUT_DIR = os.path.join(os.path.dirname(__file__), 'presentation_finale')
os.makedirs(OUT_DIR, exist_ok=True)

# Nguyen template palette
DEEP_TEAL = RGBColor(0x0D, 0x4F, 0x4F)
CREAM = RGBColor(0xFA, 0xF8, 0xF0)
GOLD = RGBColor(0xD4, 0xA8, 0x43)
TEAL = RGBColor(0x0E, 0x7C, 0x7B)
TERRACOTTA = RGBColor(0xE0, 0x7A, 0x5F)
DANGER_RED = RGBColor(0xC0, 0x39, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x2E)
DETAIL_GRAY = RGBColor(0x4A, 0x4A, 0x5A)
SUBTITLE_GRAY = RGBColor(0x8A, 0x8A, 0x9A)
LIGHT_TEAL_BG = RGBColor(0xE0, 0xF2, 0xF1)
LIGHT_ORANGE_BG = RGBColor(0xFF, 0xF3, 0xE0)
ZEBRA_GRAY = RGBColor(0xF7, 0xF7, 0xF7)
TITLE_SUBTITLE = RGBColor(0xB0, 0xD4, 0xD4)
TITLE_META = RGBColor(0x7F, 0xB8, 0xB8)

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

FONT_TITLE = 'Georgia'
FONT_BODY = 'Calibri'


# ── Helpers ──────────────────────────────────────────────────

def add_cream_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = CREAM
    bg.line.fill.background()


def add_slide_title(slide, title, subtitle=None):
    """Georgia bold title + optional italic subtitle (Nguyen pattern)."""
    tb = slide.shapes.add_textbox(Inches(0.60), Inches(0.25), Inches(8.5), Inches(0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT_TITLE
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = DEEP_TEAL
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.60), Inches(0.72), Inches(8.5), Inches(0.3))
        tf2 = tb2.text_frame
        tf2.margin_left = tf2.margin_right = tf2.margin_top = tf2.margin_bottom = 0
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.name = FONT_BODY
        run2.font.size = Pt(12)
        run2.font.italic = True
        run2.font.color.rgb = SUBTITLE_GRAY


def add_text(slide, left, top, width, height, text, size=12, bold=False,
             color=DETAIL_GRAY, align=PP_ALIGN.LEFT, font=None, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(91440)
    tf.margin_right = Emu(91440)
    tf.margin_top = Emu(45720)
    tf.margin_bottom = Emu(45720)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font or FONT_BODY
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    p.alignment = align
    return tf


def add_multiline(slide, left, top, width, height, lines, size=11,
                   color=DETAIL_GRAY, bold=False, spacing=None):
    """Multiple lines with per-line control."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(91440)
    tf.margin_right = Emu(91440)
    tf.margin_top = Emu(45720)
    tf.margin_bottom = Emu(45720)
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if spacing:
            p.space_after = Pt(spacing)
        if isinstance(line, tuple):
            txt, lcolor, lbold = line
        else:
            txt, lcolor, lbold = line, color, bold
        run = p.add_run()
        run.text = txt
        run.font.name = FONT_BODY
        run.font.size = Pt(size)
        run.font.bold = lbold
        run.font.color.rgb = lcolor
    return tf


def add_card(slide, left, top, width, height, accent_color):
    """White card with thin colored top bar (Nguyen pattern)."""
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  left, top, width, Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    return card


def add_badge(slide, left, top, number, color):
    """Numbered colored badge (Nguyen pattern)."""
    badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    left, top, Inches(0.40), Inches(0.40))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(number)
    run.font.name = FONT_BODY
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def add_insight_bar(slide, top, text, bg_color=LIGHT_TEAL_BG, text_color=DEEP_TEAL):
    """Bottom insight/callout bar (Nguyen pattern)."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.40), top, Inches(9.20), Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = bg_color
    bar.line.fill.background()
    add_text(slide, Inches(0.55), top, Inches(9.00), Inches(0.55),
             text, 12, True, text_color, PP_ALIGN.LEFT)


def add_table(slide, left, top, width, height, data, col_widths=None):
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = ts.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = FONT_BODY
                    run.font.size = Pt(10)
                    if r == 0:
                        run.font.bold = True
                        run.font.color.rgb = WHITE
                    else:
                        run.font.color.rgb = NEAR_BLACK
                para.alignment = PP_ALIGN.CENTER
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DEEP_TEAL
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ZEBRA_GRAY
    return table


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
bg.fill.solid()
bg.fill.fore_color.rgb = DEEP_TEAL
bg.line.fill.background()

# Gold vertical accent bar
vbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0.60), Inches(0.90), Inches(0.06), Inches(1.80))
vbar.fill.solid()
vbar.fill.fore_color.rgb = GOLD
vbar.line.fill.background()

# Title
add_text(slide, Inches(0.90), Inches(0.80), Inches(8.0), Inches(1.0),
         "HPL Benchmark\nImplémentation et résultats GPU", 30, True, WHITE,
         font=FONT_TITLE)

# Subtitle
add_text(slide, Inches(0.90), Inches(2.0), Inches(8.0), Inches(0.8),
         "Analyse comparative A100 (Ampere) vs H100 (Hopper)\n"
         "Partie 2 : Exécution, Résultats, Analyse", 16, False, TITLE_SUBTITLE)

# Gold horizontal divider
hbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               0, Inches(3.90), SLIDE_W, Inches(0.03))
hbar.fill.solid()
hbar.fill.fore_color.rgb = GOLD
hbar.line.fill.background()

# Presenter info
add_text(slide, Inches(0.90), Inches(4.05), Inches(8.0), Inches(0.3),
         "BENMALK Achraf", 14, True, GOLD)
add_text(slide, Inches(0.90), Inches(4.35), Inches(8.0), Inches(0.3),
         "Projet MPNA — Méthodes et Programmation Numérique Avancée — 2026", 11, False, TITLE_META)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Plateforme d'exécution
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_cream_bg(slide)
add_slide_title(slide, "Plateforme d'exécution",
                "Cluster HPC avec GPUs NVIDIA — Conteneur HPC-Benchmarks 23.10")

# Left card: Environment
add_card(slide, Inches(0.40), Inches(1.20), Inches(4.20), Inches(2.80), TEAL)
add_text(slide, Inches(0.50), Inches(1.35), Inches(4.0), Inches(0.35),
         "Environnement", 13, True, TEAL)
add_multiline(slide, Inches(0.50), Inches(1.75), Inches(4.0), Inches(2.0), [
    "Conteneur NVIDIA HPC-Benchmarks 23.10",
    "Déploiement via Singularity",
    "Bibliothèques : cuBLAS + NCCL",
    "Ordonnanceur : Slurm",
    "Tensor Cores FP64 pour DGEMM",
], 11, DETAIL_GRAY, spacing=4)

# Right card: GPU comparison table
add_card(slide, Inches(4.80), Inches(1.20), Inches(5.00), Inches(2.80), GOLD)
add_text(slide, Inches(4.90), Inches(1.35), Inches(4.8), Inches(0.35),
         "GPUs testés", 13, True, GOLD)
gpu_data = [
    ["", "A100 (Ampere)", "H100 (Hopper)"],
    ["Partition", "gpu", "gpu_h100"],
    ["Mémoire", "80 Go HBM2e", "80 Go HBM3"],
    ["Cœurs CUDA", "6 912", "16 896"],
    ["Peak FP64", "19,5 TFLOPS", "54 TFLOPS"],
]
add_table(slide, Inches(4.90), Inches(1.80), Inches(4.80), Inches(2.0), gpu_data)

# Bottom insight bar
add_insight_bar(slide, Inches(4.25),
    "Configuration HPL : NB = 576 (blocs larges pour GPU)  |  P × Q = 2 × 1  |  BCAST = 6")

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Étapes d'exécution (4 numbered cards, 2x2)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_cream_bg(slide)
add_slide_title(slide, "Étapes d'exécution sur GPU",
                "Du fichier de configuration au lancement du benchmark")

steps = [
    (1, "HPL.dat", "Création du fichier de configuration\n(N, NB, P×Q, BCAST)", TEAL),
    (2, "Singularity", "Lancement du conteneur NVIDIA\navec bind du répertoire de travail", TEAL),
    (3, "Slurm", "Allocation GPU via salloc\n(--gres=gpu:1 ou gpu:2)", TEAL),
    (4, "Exécution", "mpirun -np <X> ./hpl.sh\n--dat /mnt/HPL.dat", TERRACOTTA),
]

card_positions = [
    (Inches(0.40), Inches(1.15)),
    (Inches(5.10), Inches(1.15)),
    (Inches(0.40), Inches(3.10)),
    (Inches(5.10), Inches(3.10)),
]

img_files = ["Step1.png", "Step2.png", "Step3.png", "Step4.png"]

for i, (num, title, desc, color) in enumerate(steps):
    cx, cy = card_positions[i]
    cw, ch = Inches(4.50), Inches(1.70)
    add_card(slide, cx, cy, cw, ch, color)
    add_badge(slide, cx + Inches(0.15), cy + Inches(0.20), num, color)
    add_text(slide, cx + Inches(0.65), cy + Inches(0.15), Inches(1.8), Inches(0.35),
             title, 13, True, color)
    add_text(slide, cx + Inches(0.65), cy + Inches(0.50), Inches(1.8), Inches(1.0),
             desc, 10, False, DETAIL_GRAY)
    img_path = os.path.join(IMG_DIR, img_files[i])
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, cx + Inches(2.6), cy + Inches(0.15), Inches(1.75))

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Résultats A100
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_cream_bg(slide)
add_slide_title(slide, "Résultats : GPU A100 (Ampere)",
                "Peak théorique : 19,5 TFLOPS  |  Configuration : NB=576, P×Q=2×1")

# Left: table
a100_data = [
    ["N", "1 GPU", "2 GPUs", "Speedup"],
    ["20 000", "9 731", "9 106", "0,94x"],
    ["40 000", "15 890", "25 230", "1,59x"],
    ["60 000", "17 150", "31 430", "1,83x"],
    ["80 000", "17 600", "33 560", "1,91x"],
    ["100 000", "17 860", "34 720", "1,95x"],
]
add_table(slide, Inches(0.40), Inches(1.10), Inches(4.00), Inches(2.5), a100_data)

# Right: graph
img_path = os.path.join(IMG_DIR, 'hpl-a100.png')
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(4.60), Inches(1.10), Inches(5.20))

# Bottom metrics cards
add_card(slide, Inches(0.40), Inches(3.85), Inches(4.30), Inches(0.80), TEAL)
add_text(slide, Inches(0.50), Inches(3.95), Inches(4.1), Inches(0.60),
         "Efficacité : 17 860 / 19 487 = 91,7%\n"
         "Speedup 2 GPUs : 1,95x → eff. parallèle 97%", 11, True, TEAL)

add_card(slide, Inches(4.90), Inches(3.85), Inches(4.90), Inches(0.80), TERRACOTTA)
add_text(slide, Inches(5.00), Inches(3.95), Inches(4.7), Inches(0.60),
         "Anomalie N=20K : 2 GPUs plus lents (-6,4%)\n"
         "Tous les tests : résidu PASSED", 11, True, TERRACOTTA)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — Résultats H100
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_cream_bg(slide)
add_slide_title(slide, "Résultats : GPU H100 (Hopper)",
                "Peak théorique : 54 TFLOPS  |  Configuration : NB=576, P×Q=2×1")

h100_data = [
    ["N", "1 GPU", "2 GPUs", "Speedup"],
    ["20 000", "16 130", "11 510", "0,71x"],
    ["40 000", "35 760", "41 460", "1,16x"],
    ["60 000", "41 760", "64 700", "1,55x"],
    ["80 000", "44 130", "76 730", "1,74x"],
    ["100 000", "45 110", "81 970", "1,82x"],
]
add_table(slide, Inches(0.40), Inches(1.10), Inches(4.00), Inches(2.5), h100_data)

img_path = os.path.join(IMG_DIR, 'hpl-h100.png')
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(4.60), Inches(1.10), Inches(5.20))

add_card(slide, Inches(0.40), Inches(3.85), Inches(4.30), Inches(0.80), TEAL)
add_text(slide, Inches(0.50), Inches(3.95), Inches(4.1), Inches(0.60),
         "Efficacité : 45 110 / 54 067 = 83,4%\n"
         "Speedup 2 GPUs : 1,82x → eff. parallèle 91%", 11, True, TEAL)

add_card(slide, Inches(4.90), Inches(3.85), Inches(4.90), Inches(0.80), DANGER_RED)
add_text(slide, Inches(5.00), Inches(3.95), Inches(4.7), Inches(0.60),
         "Anomalie N=20K : 2 GPUs 28,6% PLUS LENTS !\n"
         "Tous les tests : résidu PASSED", 11, True, DANGER_RED)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — Anomalie multi-GPU
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_cream_bg(slide)
add_slide_title(slide, "Anomalie : 2 GPUs plus lents qu'un seul",
                "Dégradation de performance observée à N = 20 000")

# Left: anomaly table
anom_data = [
    ["GPU", "1 GPU", "2 GPUs", "Dégradation"],
    ["A100", "9 731", "9 106", "-6,4%"],
    ["H100", "16 130", "11 510", "-28,6%"],
]
add_table(slide, Inches(0.40), Inches(1.10), Inches(4.20), Inches(1.0), anom_data)

# Left: numbered explanation cards
explanations = [
    (1, "Problème trop petit — chaque GPU reçoit une portion insuffisante", TEAL),
    (2, "Coût de communication inter-GPU > gain de calcul", TEAL),
    (3, "Plus prononcé sur H100 : 16 896 vs 6 912 cœurs", TEAL),
    (4, "Seuil de rentabilité : entre N = 20K et 40K", TERRACOTTA),
]
for i, (num, txt, color) in enumerate(explanations):
    y = Inches(2.30 + i * 0.55)
    add_badge(slide, Inches(0.50), y, num, color)
    add_text(slide, Inches(1.00), y, Inches(3.5), Inches(0.45),
             txt, 11, False, NEAR_BLACK)

# Right: explanation card
add_card(slide, Inches(4.80), Inches(1.10), Inches(5.00), Inches(3.30), GOLD)
add_text(slide, Inches(4.90), Inches(1.25), Inches(4.8), Inches(0.30),
         "Pourquoi ?", 13, True, GOLD)
add_multiline(slide, Inches(4.90), Inches(1.65), Inches(4.8), Inches(2.6), [
    ("Le calcul croît en O(N³)", DEEP_TEAL, True),
    "Volume de travail augmente cubiquement avec N",
    "",
    ("La communication croît en O(N²)", TERRACOTTA, True),
    "Échanges inter-GPU augmentent quadratiquement",
    "",
    ("Pour les grands N, le calcul domine", DEEP_TEAL, True),
    "Les GPUs passent plus de temps à calculer qu'à communiquer",
    "",
    ("Pour les petits N, la communication domine", TERRACOTTA, True),
    "Le surcoût de synchronisation dépasse le bénéfice du parallélisme",
], 10, DETAIL_GRAY, spacing=2)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Comparaison A100 vs H100
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_cream_bg(slide)
add_slide_title(slide, "Comparaison architecturale : A100 vs H100",
                "Facteur d'accélération en fonction de la taille du problème")

# Left: graph
img_path = os.path.join(IMG_DIR, 'a100-h100-hpl.png')
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.30), Inches(1.05), Inches(4.80))

# Right: ratio table
ratio_data = [
    ["N", "A100", "H100", "Ratio"],
    ["20K", "9 731", "16 130", "1,66x"],
    ["40K", "15 890", "35 760", "2,25x"],
    ["60K", "17 150", "41 760", "2,43x"],
    ["80K", "17 600", "44 130", "2,51x"],
    ["100K", "17 860", "45 110", "2,53x"],
]
add_table(slide, Inches(5.30), Inches(1.05), Inches(4.50), Inches(2.5), ratio_data)

# Right: ratio analysis card
add_card(slide, Inches(5.30), Inches(3.70), Inches(4.50), Inches(1.20), TEAL)
add_multiline(slide, Inches(5.40), Inches(3.80), Inches(4.3), Inches(1.0), [
    ("Ratio théorique : 54 / 19,5 = 2,77x", TEAL, True),
    ("Ratio mesuré (N=100K) : 2,53x", DEEP_TEAL, True),
    ("Écart ~9% — dû à la différence d'efficacité", DETAIL_GRAY, False),
    ("H100 : 83,4%  vs  A100 : 91,7%", DETAIL_GRAY, False),
], 11, spacing=3)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — Synthèse de l'efficacité
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_cream_bg(slide)
add_slide_title(slide, "Synthèse de l'efficacité",
                "Efficacité = GFLOPS mesurés / Pic théorique × 100%")

eff_data = [
    ["Configuration", "Peak (TFLOPS)", "Meilleur (GFLOPS)", "Efficacité", "Eff. parallèle"],
    ["A100 (1 GPU)", "19,5", "17 860", "91,7%", "—"],
    ["H100 (1 GPU)", "54,0", "45 110", "83,4%", "—"],
    ["A100 (2 GPUs)", "39,0", "34 720", "89,0%", "97%"],
    ["H100 (2 GPUs)", "108,0", "81 970", "75,8%", "91%"],
]
add_table(slide, Inches(0.40), Inches(1.05), Inches(9.20), Inches(2.2), eff_data)

# Two insight cards
add_card(slide, Inches(0.40), Inches(3.50), Inches(4.40), Inches(1.10), TEAL)
add_multiline(slide, Inches(0.50), Inches(3.60), Inches(4.2), Inches(0.90), [
    ("HPL est compute-bound", TEAL, True),
    ("Le calcul (DGEMM) domine le temps d'exécution.", DETAIL_GRAY, False),
    ("L'A100 est plus facile à saturer → meilleure efficacité.", DETAIL_GRAY, False),
], 10, spacing=2)

add_card(slide, Inches(5.00), Inches(3.50), Inches(4.80), Inches(1.10), TERRACOTTA)
add_multiline(slide, Inches(5.10), Inches(3.60), Inches(4.6), Inches(0.90), [
    ("Plus de puissance = plus difficile à exploiter", TERRACOTTA, True),
    ("L'efficacité diminue de 91,7% (A100 1 GPU)", DETAIL_GRAY, False),
    ("à 75,8% (H100 2 GPUs).", DETAIL_GRAY, False),
], 10, spacing=2)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — Conclusion
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_cream_bg(slide)
add_slide_title(slide, "Conclusion et enseignements",
                "5 résultats clés de notre analyse HPL sur GPU")

conclusions = [
    (1, "HPL est compute-bound",
     "Efficacités de 83 à 92% — le calcul domine", TEAL),
    (2, "La taille du problème est déterminante",
     "O(N³) calcul vs O(N²) communication", TEAL),
    (3, "Le multi-GPU a un seuil de rentabilité",
     "N entre 20K et 40K — en dessous, contre-productif", TEAL),
    (4, "Le H100 offre ≈2,5x l'A100",
     "Ratio mesuré 2,53x vs théorique 2,77x", TERRACOTTA),
    (5, "L'efficacité diminue avec la puissance",
     "97% (A100) vs 91% (H100) en parallèle", TERRACOTTA),
]

for i, (num, title, desc, color) in enumerate(conclusions):
    y = Inches(1.05 + i * 0.75)
    add_card(slide, Inches(0.40), y, Inches(9.20), Inches(0.65), color)
    add_badge(slide, Inches(0.55), y + Inches(0.12), num, color)
    add_text(slide, Inches(1.10), y + Inches(0.05), Inches(4.0), Inches(0.30),
             title, 12, True, color)
    add_text(slide, Inches(1.10), y + Inches(0.33), Inches(8.3), Inches(0.25),
             desc, 10, False, DETAIL_GRAY)

add_insight_bar(slide, Inches(4.95),
    "Puissance brute, passage à l'échelle et dimensionnement du problème sont étroitement liés.")

# ── Save ─────────────────────────────────────────────────────
output_path = os.path.join(OUT_DIR, 'HPL_Resultats_Analyse.pptx')
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
