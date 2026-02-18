#!/usr/bin/env python3
"""Assemble HPL2.pptx (Part 1: code analysis) + HPL.pptx (Part 2: GPU results)
into a single presentation with a Plan slide at the beginning.
NO existing slides are modified — only a Plan/sommaire slide is added."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy
from lxml import etree
import os

BASE_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'presentation_finale')

# Load source presentations
prs_part1 = Presentation(os.path.join(BASE_DIR, 'HPL2.pptx'))   # Binome: code analysis
prs_part2 = Presentation(os.path.join(BASE_DIR, 'HPL.pptx'))    # User: GPU results

# Create destination — match 16:9 dimensions
dst = Presentation()
dst.slide_width = prs_part2.slide_width    # 10" x 5.625"
dst.slide_height = prs_part2.slide_height

# ── Palette (Nguyen style) ─────────────────────────────
DEEP_TEAL   = RGBColor(0x0D, 0x4F, 0x4F)
CREAM       = RGBColor(0xFA, 0xF8, 0xF0)
GOLD        = RGBColor(0xD4, 0xA8, 0x43)
TEAL        = RGBColor(0x0E, 0x7C, 0x7B)
TERRACOTTA  = RGBColor(0xE0, 0x7A, 0x5F)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK  = RGBColor(0x1A, 0x1A, 0x2E)
DETAIL_GRAY = RGBColor(0x4A, 0x4A, 0x5A)
SUBTITLE_GRAY = RGBColor(0x8A, 0x8A, 0x9A)
TITLE_SUBTITLE = RGBColor(0xB0, 0xD4, 0xD4)
TITLE_META  = RGBColor(0x7F, 0xB8, 0xB8)

FONT_TITLE = 'Georgia'
FONT_BODY  = 'Calibri'

R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'


# ── Slide copy helper ──────────────────────────────────

def copy_slide_from(src_prs, src_index, dst_prs):
    """Copy a slide from src_prs into dst_prs, preserving shapes, images, tables."""
    src_slide = src_prs.slides[src_index]
    layout = dst_prs.slide_layouts[6]  # Blank
    new_slide = dst_prs.slides.add_slide(layout)

    # 1) Collect non-layout, non-notes relationships to copy
    rels_to_copy = {}
    for key, rel in src_slide.part.rels.items():
        if rel.is_external:
            continue
        rt = rel.reltype
        if 'slideLayout' in rt or 'notesSlide' in rt or 'theme' in rt:
            continue
        rels_to_copy[key] = rel

    # 2) Copy each relationship → build rId mapping (old → new)
    rId_map = {}
    for old_rId, rel in rels_to_copy.items():
        new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
        rId_map[old_rId] = new_rId

    # 3) Deep-copy source shapes into destination spTree
    src_spTree = src_slide.shapes._spTree
    dst_spTree = new_slide.shapes._spTree

    keep_tags = {'nvGrpSpPr', 'grpSpPr'}
    for child in list(dst_spTree):
        if etree.QName(child).localname not in keep_tags:
            dst_spTree.remove(child)

    for child in src_spTree:
        if etree.QName(child).localname not in keep_tags:
            dst_spTree.append(deepcopy(child))

    # 4) Remap rId references in the copied XML
    for old_rId, new_rId in rId_map.items():
        if old_rId != new_rId:
            for elem in dst_spTree.iter():
                for attr in [f'{{{R_NS}}}embed', f'{{{R_NS}}}link', f'{{{R_NS}}}id']:
                    if elem.get(attr) == old_rId:
                        elem.set(attr, new_rId)

    # 5) Copy speaker notes if present
    try:
        if src_slide.has_notes_slide:
            src_notes = src_slide.notes_slide
            dst_notes = new_slide.notes_slide
            dst_tf = dst_notes.notes_text_frame
            src_tf = src_notes.notes_text_frame
            dst_tf.text = src_tf.text
    except Exception:
        pass

    return new_slide


# ═══════════════════════════════════════════════════════
# SLIDE 1 — Title slide
# ═══════════════════════════════════════════════════════
slide = dst.slides.add_slide(dst.slide_layouts[6])

# Dark teal background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, dst.slide_width, dst.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DEEP_TEAL
bg.line.fill.background()

# Gold vertical accent bar
vbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0.60), Inches(0.90), Inches(0.06), Inches(1.80))
vbar.fill.solid()
vbar.fill.fore_color.rgb = GOLD
vbar.line.fill.background()

# Title
tb = slide.shapes.add_textbox(Inches(0.90), Inches(0.80), Inches(8.0), Inches(1.0))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "HPL Benchmark\nImplémentation et résultats GPU"
run.font.name = FONT_TITLE
run.font.size = Pt(30)
run.font.bold = True
run.font.color.rgb = WHITE

# Subtitle
tb = slide.shapes.add_textbox(Inches(0.90), Inches(2.0), Inches(8.0), Inches(0.8))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Analyse du code source et résultats sur GPU\nA100 (Ampere) vs H100 (Hopper)"
run.font.name = FONT_BODY
run.font.size = Pt(16)
run.font.color.rgb = TITLE_SUBTITLE

# Gold horizontal divider
hbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               0, Inches(3.90), dst.slide_width, Inches(0.03))
hbar.fill.solid()
hbar.fill.fore_color.rgb = GOLD
hbar.line.fill.background()

# Authors
tb = slide.shapes.add_textbox(Inches(0.90), Inches(4.05), Inches(8.0), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "BENMALK Achraf  —  KARIMI Karim"
run.font.name = FONT_BODY
run.font.size = Pt(14)
run.font.bold = True
run.font.color.rgb = GOLD

tb = slide.shapes.add_textbox(Inches(0.90), Inches(4.35), Inches(8.0), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Projet MPNA — Méthodes et Programmation Numérique Avancée — 2026"
run.font.name = FONT_BODY
run.font.size = Pt(11)
run.font.color.rgb = TITLE_META


# ═══════════════════════════════════════════════════════
# SLIDE 2 — Plan / Sommaire
# ═══════════════════════════════════════════════════════
slide = dst.slides.add_slide(dst.slide_layouts[6])

# Cream background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, dst.slide_width, dst.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = CREAM
bg.line.fill.background()

# Title
tb = slide.shapes.add_textbox(Inches(0.60), Inches(0.25), Inches(8.5), Inches(0.5))
tf = tb.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Plan de la présentation"
run.font.name = FONT_TITLE
run.font.size = Pt(24)
run.font.bold = True
run.font.color.rgb = DEEP_TEAL

# ── Part 1 card ──
card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.50), Inches(1.05), Inches(9.00), Inches(1.90))
card.fill.solid()
card.fill.fore_color.rgb = WHITE
card.line.fill.background()

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.50), Inches(1.05), Inches(9.00), Inches(0.05))
bar.fill.solid()
bar.fill.fore_color.rgb = TEAL
bar.line.fill.background()

# Badge 1
badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.70), Inches(1.25), Inches(0.40), Inches(0.40))
badge.fill.solid()
badge.fill.fore_color.rgb = TEAL
badge.line.fill.background()
btf = badge.text_frame
btf.vertical_anchor = MSO_ANCHOR.MIDDLE
btf.margin_left = btf.margin_right = btf.margin_top = btf.margin_bottom = 0
p = btf.paragraphs[0]
run = p.add_run()
run.text = "1"
run.font.name = FONT_BODY
run.font.size = Pt(14)
run.font.bold = True
run.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Part 1 title
tb = slide.shapes.add_textbox(Inches(1.30), Inches(1.20), Inches(7.5), Inches(0.40))
tf = tb.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Partie 1 — Analyse du code source (Mini HPL)"
run.font.name = FONT_TITLE
run.font.size = Pt(14)
run.font.bold = True
run.font.color.rgb = TEAL

# Part 1 items
tb = slide.shapes.add_textbox(Inches(1.30), Inches(1.65), Inches(8.0), Inches(1.2))
tf = tb.text_frame
tf.word_wrap = True
items_p1 = [
    "Élimination de Gauss et pivotage partiel",
    "Exemple pas à pas (tournage à la main)",
    "Communications MPI : recherche du pivot, échange de lignes",
    "Broadcast, mise à jour, back-substitution parallèle",
    "Calcul du résidu et validation",
]
for i, item in enumerate(items_p1):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = f"  •  {item}"
    run.font.name = FONT_BODY
    run.font.size = Pt(11)
    run.font.color.rgb = DETAIL_GRAY
    p.space_after = Pt(2)

# ── Part 2 card ──
card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.50), Inches(3.20), Inches(9.00), Inches(1.90))
card.fill.solid()
card.fill.fore_color.rgb = WHITE
card.line.fill.background()

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.50), Inches(3.20), Inches(9.00), Inches(0.05))
bar.fill.solid()
bar.fill.fore_color.rgb = TERRACOTTA
bar.line.fill.background()

# Badge 2
badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.70), Inches(3.40), Inches(0.40), Inches(0.40))
badge.fill.solid()
badge.fill.fore_color.rgb = TERRACOTTA
badge.line.fill.background()
btf = badge.text_frame
btf.vertical_anchor = MSO_ANCHOR.MIDDLE
btf.margin_left = btf.margin_right = btf.margin_top = btf.margin_bottom = 0
p = btf.paragraphs[0]
run = p.add_run()
run.text = "2"
run.font.name = FONT_BODY
run.font.size = Pt(14)
run.font.bold = True
run.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Part 2 title
tb = slide.shapes.add_textbox(Inches(1.30), Inches(3.35), Inches(7.5), Inches(0.40))
tf = tb.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Partie 2 — Exécution et résultats GPU"
run.font.name = FONT_TITLE
run.font.size = Pt(14)
run.font.bold = True
run.font.color.rgb = TERRACOTTA

# Part 2 items
tb = slide.shapes.add_textbox(Inches(1.30), Inches(3.80), Inches(8.0), Inches(1.2))
tf = tb.text_frame
tf.word_wrap = True
items_p2 = [
    "Plateforme d'exécution et étapes du benchmark",
    "Résultats A100 (Ampere) et H100 (Hopper)",
    "Anomalie multi-GPU à N = 20 000",
    "Comparaison architecturale et synthèse de l'efficacité",
    "Conclusion et enseignements",
]
for i, item in enumerate(items_p2):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = f"  •  {item}"
    run.font.name = FONT_BODY
    run.font.size = Pt(11)
    run.font.color.rgb = DETAIL_GRAY
    p.space_after = Pt(2)


# ═══════════════════════════════════════════════════════
# COPY Part 1 slides (HPL2.pptx — binome's code analysis)
# ═══════════════════════════════════════════════════════
print("Copying Part 1 (HPL2.pptx)...")
for i in range(len(prs_part1.slides)):
    print(f"  Slide {i+1}/{len(prs_part1.slides)}")
    copy_slide_from(prs_part1, i, dst)

# ═══════════════════════════════════════════════════════
# COPY Part 2 slides (HPL.pptx — user's GPU results)
# ═══════════════════════════════════════════════════════
print("Copying Part 2 (HPL.pptx)...")
for i in range(len(prs_part2.slides)):
    print(f"  Slide {i+1}/{len(prs_part2.slides)}")
    copy_slide_from(prs_part2, i, dst)

# ── Save ──────────────────────────────────────────────
output_path = os.path.join(BASE_DIR, 'HPL_Final.pptx')
dst.save(output_path)
print(f"\nAssembled presentation saved to: {output_path}")
print(f"Total slides: {len(dst.slides)}")
print(f"  - 1 Title slide (new)")
print(f"  - 1 Plan slide (new)")
print(f"  - {len(prs_part1.slides)} Part 1 slides (HPL2.pptx)")
print(f"  - {len(prs_part2.slides)} Part 2 slides (HPL.pptx)")
