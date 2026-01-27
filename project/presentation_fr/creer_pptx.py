#!/usr/bin/env python3
"""
Génère le fichier PowerPoint pour la présentation HPL
Partie 2: Slides 9-14 (Résultats)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Créer la présentation
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

def rgb_color(r, g, b):
    """Helper to create RGB color"""
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)

# Couleurs
DARK_BLUE = rgb_color(44, 62, 80)
GREEN = rgb_color(39, 174, 96)
ORANGE = rgb_color(230, 126, 34)
WHITE = rgb_color(255, 255, 255)
LIGHT_GRAY = rgb_color(236, 240, 241)

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, content_lines, has_image=False, image_path=None):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BLUE
    title_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Content
    if has_image and image_path and os.path.exists(image_path):
        # Image on right
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5.5))
        slide.shapes.add_picture(image_path, Inches(6.8), Inches(1.5), width=Inches(6))
    else:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))

    tf = content_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle formatting
        if line.startswith("• "):
            p.text = line
        elif line.startswith("[BOLD]"):
            p.text = line[6:]
            p.font.bold = True
        elif line.startswith("[GREEN]"):
            p.text = line[7:]
            p.font.color.rgb = GREEN
            p.font.bold = True
        elif line.startswith("[ORANGE]"):
            p.text = line[8:]
            p.font.color.rgb = ORANGE
        elif line == "":
            p.text = ""
        else:
            p.text = line

        p.font.size = Pt(18)
        p.space_after = Pt(6)

    return slide

# ============================================
# SLIDE 9: Configuration Expérimentale
# ============================================
slide9 = add_content_slide(prs, "Configuration Expérimentale", [
    "[BOLD]Plateforme de Test",
    "• Machine : ASUS Zephyrus G14 (2020)",
    "• CPU : AMD Ryzen 9 4900HS (8 coeurs / 16 threads)",
    "• RAM : 32 Go",
    "• OS : Ubuntu 22.04 via WSL2",
    "",
    "[BOLD]Logiciels",
    "• HPL version 2.3",
    "• OpenMPI (communication inter-processus)",
    "• OpenBLAS (bibliotheque BLAS)",
    "",
    "[BOLD]Parametres HPL",
    "• NB = 128 et 192 (tailles de bloc testees)",
    "• P x Q = 2 x 4 (grille de 8 processus)",
])

# ============================================
# SLIDE 10: Résultats des Expériences
# ============================================
slide10 = add_content_slide(prs, "Resultats des Experiences", [
    "[BOLD]Tableau des Resultats",
    "",
    "• N=10 000, NB=192 : 15.7 GFLOPS (42s) - VALIDE",
    "• N=20 000, NB=192 : 41.0 GFLOPS (2min 10s) - VALIDE",
    "• N=30 000, NB=192 : 38.2 GFLOPS (7min 51s) - VALIDE",
    "[GREEN]• N=30 000, NB=128 : 44.2 GFLOPS (6min 48s) - MEILLEUR!",
    "",
    "[BOLD]Observations Cles",
    "[GREEN]Tous les tests VALIDES (PASSED)",
    "• GFLOPS augmente avec N (meilleur ratio calcul/comm.)",
    "[ORANGE]Baisse a N=30K/NB=192 due au thermal throttling",
    "[GREEN]Tuning NB: +15.6% avec NB=128 vs NB=192",
], has_image=True, image_path="/home/user/mpna-local/project/presentation_fr/graphique4_tableau.png")

# ============================================
# SLIDE 11: Évolution de la Performance
# ============================================
slide11 = add_content_slide(prs, "Evolution de la Performance", [
    "[BOLD]Tendance Observee",
    "",
    "• N=10K : 15.7 GFLOPS (reference)",
    "• N=20K : 41.0 GFLOPS (+161%)",
    "[ORANGE]• N=30K, NB=192 : 38.2 GFLOPS (throttling)",
    "[GREEN]• N=30K, NB=128 : 44.2 GFLOPS (optimise)",
    "",
    "[BOLD]Explication",
    "• Plus N est grand, meilleur ratio calcul/comm.",
    "• Le tuning de NB recupere +15.6% de perf.",
    "• Thermal throttling apres ~5min de calcul",
], has_image=True, image_path="/home/user/mpna-local/project/presentation_fr/graphique3_evolution.png")

# ============================================
# SLIDE 12: Analyse de l'Efficacité
# ============================================
slide12 = add_content_slide(prs, "Analyse de l'Efficacite", [
    "[BOLD]Calcul de l'Efficacite",
    "• Efficacite = GFLOPS obtenus / GFLOPS theoriques x 100%",
    "• Pic theorique Ryzen 9: ~400 GFLOPS (estimation)",
    "[GREEN]Meilleur resultat: 44.2 GFLOPS = Efficacite ~11%",
    "",
    "[BOLD]Pourquoi seulement 11% ?",
    "[ORANGE]• WSL2: overhead de virtualisation",
    "[ORANGE]• OpenBLAS: moins optimise qu'Intel MKL",
    "[ORANGE]• Laptop: limites thermiques vs serveur HPC",
    "",
    "[BOLD]Reference",
    "• Sur cluster HPC reel: 70-85% d'efficacite attendue",
], has_image=True, image_path="/home/user/mpna-local/project/presentation_fr/graphique5_efficacite.png")

# ============================================
# SLIDE 13: Observations et Limites
# ============================================
slide13 = add_content_slide(prs, "Observations et Limites", [
    "[BOLD]Ce Qui Fonctionne",
    "[GREEN]• Performance scale avec N jusqu'aux limites thermiques",
    "[GREEN]• Tous les tests passent la validation numerique",
    "[GREEN]• Resultats reproductibles et coherents",
    "[GREEN]• Le tuning NB ameliore significativement les perfs",
    "",
    "[BOLD]Limites de Notre Setup",
    "[ORANGE]• WSL2: couche de virtualisation = overhead",
    "[ORANGE]• Laptop: refroidissement limite = throttling",
    "[ORANGE]• OpenBLAS: pas optimise pour AMD Ryzen",
    "",
    "[BOLD]Pour Ameliorer",
    "• Cluster HPC dedie avec refroidissement adapte",
    "• Intel MKL ou AMD BLIS (bibliotheques optimisees)",
])

# ============================================
# SLIDE 14: Conclusion
# ============================================
slide14 = add_content_slide(prs, "Conclusion", [
    "[BOLD]Ce Qu'on Retient",
    "",
    "[GREEN]HPL mesure les GFLOPS via resolution de Ax = b",
    "[GREEN]Algorithme: Decomposition LU avec pivotage partiel",
    "[GREEN]Parallelisation: Distribution 2D block-cyclic",
    "[GREEN]Nos resultats: Valides, coherents avec la theorie",
    "",
    "[BOLD]Meilleur Resultat",
    "[GREEN]44.2 GFLOPS avec N=30 000 et NB=128",
    "",
    "[BOLD]Tendance Cle",
    "• GFLOPS augmente avec N (meilleur ratio calcul/comm.)",
    "• Le tuning des parametres est important (+15.6%)",
])

# ============================================
# SLIDE FINAL: Questions
# ============================================
add_title_slide(prs, "Merci !", "Questions ?")

# Sauvegarder
output_path = "/home/user/mpna-local/project/presentation_fr/HPL_Presentation_Resultats.pptx"
prs.save(output_path)
print(f"Presentation creee: {output_path}")
